# -*- coding: utf-8 -*-

import json
from pptx import Presentation

def get_text(filepath):
    ppt = Presentation(filepath) # Read the Slides
    slides_text = []
    title = [] # Title of each slide
    content = [] # The entire content of ppt
    slides_index = 1
    for slide in ppt.slides:
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text:
                slide_text.append(shape.text_frame.text)
        title.append(slide_text[0])
        content.append(slide_text[1:])
        slides_text.append("## Slide"+str(slides_index)+''.join(slide_text)+"\n\n")
        slides_index+=1
        ppt_md = ''.join(slides_text)
    return ppt_md, title, content

def build_template(title, content):
    template_data = []
    # check length
    if len(title) != len(content):
        return("Please check the ppt slides title and content!")
    for index in range(len(title)):
        item = {
            "index":index+1,
            "title":title[index],
            "content":content[index],
            "script":""
        }
        template_data.append(item)
    json_data = json.loads(json.dumps(template_data, ensure_ascii=False, indent=4))
    return json_data
